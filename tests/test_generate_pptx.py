"""Tests for PowerPoint generation."""

import os
import tempfile

from pptx import Presentation

from scripts.generate_pptx import generate_sample_pptx


def test_generate_sample_pptx_creates_file():
    """Verify that generate_sample_pptx produces a valid .pptx file."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = generate_sample_pptx(output_dir=tmpdir)
        assert os.path.exists(path)
        assert path.endswith(".pptx")


def test_generated_pptx_has_expected_slides():
    """Verify the generated presentation has the correct number of slides."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = generate_sample_pptx(output_dir=tmpdir)
        prs = Presentation(path)
        assert len(prs.slides) == 4  # title, content, chart, table


def test_generated_pptx_is_widescreen():
    """Verify the generated presentation uses 16:9 aspect ratio."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = generate_sample_pptx(output_dir=tmpdir)
        prs = Presentation(path)
        # 16:9 = ~13.333 x 7.5 inches; allow small EMU rounding tolerance
        assert abs(prs.slide_width - 12192000) < 1000
        assert abs(prs.slide_height - 6858000) < 1000


def test_generated_pptx_title_slide():
    """Verify the title slide has correct content."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = generate_sample_pptx(output_dir=tmpdir)
        prs = Presentation(path)
        title_slide = prs.slides[0]
        assert title_slide.shapes.title.text == "Sample Presentation"
