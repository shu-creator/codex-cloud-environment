"""Generate a sample PowerPoint presentation to verify the harness works."""

import os
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def generate_sample_pptx(output_dir: str = "output") -> str:
    """Generate a sample presentation with various slide types.

    Returns the path to the generated file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = "Document Harness Verification | Auto-generated"

    # --- Content Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Highlights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Harness capabilities:"
    for point in [
        "Automated PPTX generation",
        "Chart and table support",
        "Template-based creation",
        "Datetime-stamped output",
    ]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1

    # --- Chart Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1)
    )
    title_shape.text_frame.text = "Quarterly Performance"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("2024", (120, 135, 148, 162))
    chart_data.add_series("2025", (145, 158, 172, 190))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.5),
        Inches(1.8),
        Inches(10),
        Inches(5),
        chart_data,
    ).chart
    chart.has_legend = True

    # --- Table Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(1)
    )
    title_shape.text_frame.text = "Feature Comparison"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    table = slide.shapes.add_table(
        4, 3, Inches(2), Inches(2), Inches(9), Inches(3.5)
    ).table
    headers = ["Feature", "Plan A", "Plan B"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    data = [
        ("Slides", "10", "20"),
        ("Charts", "Yes", "Yes"),
        ("Templates", "Basic", "Premium"),
    ]
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, val in enumerate(row_data):
            table.cell(row_idx, col_idx).text = val

    # Save
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(output_dir, f"presentation_{timestamp}.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_sample_pptx()
