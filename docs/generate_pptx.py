#!/usr/bin/env python3
"""3つのHTML資料をPowerPointに変換するスクリプト."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT_DIR = Path(__file__).parent

# ── 共通色定義 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
LIGHT_RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
SLIDE_BG = RGBColor(0xF3, 0xF4, 0xF6)


# ──────────── ヘルパー関数 ────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, bg_color=DARK_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, bg_color)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, section_num, section_title, bg_color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{section_num}. {section_title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, body_items, title_color=BLUE, bullet_color=BLACK, bg_color=SLIDE_BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = title_color

    # 下線
    line_shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.1), Inches(11.8), Pt(3)  # MSO_SHAPE.RECTANGLE
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = title_color
    line_shape.line.fill.background()

    # ボディ
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.4), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(body_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        if isinstance(item, tuple):
            text, font_size, bold, color, indent = item
        else:
            text = item
            font_size = Pt(16)
            bold = False
            color = bullet_color
            indent = 0

        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.level = indent
        p.space_after = Pt(4)

    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items,
                         left_color=GREEN, right_color=BLUE, bg_color=SLIDE_BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # 左カラム
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(20)
    lp.font.bold = True
    lp.font.color.rgb = left_color
    lp.space_after = Pt(8)
    for item in left_items:
        lp2 = ltf.add_paragraph()
        lp2.text = f"  {item}"
        lp2.font.size = Pt(14)
        lp2.font.color.rgb = BLACK
        lp2.space_after = Pt(3)

    # 右カラム
    right_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.3), Inches(5.8), Inches(5.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(20)
    rp.font.bold = True
    rp.font.color.rgb = right_color
    rp.space_after = Pt(8)
    for item in right_items:
        rp2 = rtf.add_paragraph()
        rp2.text = f"  {item}"
        rp2.font.size = Pt(14)
        rp2.font.color.rgb = BLACK
        rp2.space_after = Pt(3)

    return slide


def add_table_slide(prs, title, headers, rows, title_color=BLUE, bg_color=SLIDE_BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = title_color

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(11.8)
    col_width = table_width // num_cols

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.6), Inches(1.2), table_width, Inches(0.4 * num_rows)
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = BLACK
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

    return slide


def add_flow_slide(prs, title, steps, colors=None, bg_color=SLIDE_BG):
    """縦方向のフロー図スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color)

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(11.8), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE

    if colors is None:
        colors = [BLUE] * len(steps)

    box_width = Inches(8)
    box_height = Inches(0.9)
    x_start = Inches(3)
    y_start = Inches(1.4)
    gap = Inches(0.15)

    for i, (step_title, step_desc) in enumerate(steps):
        y = y_start + i * (box_height + Inches(0.5) + gap)

        # ボックス
        shape = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            x_start, y, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i] if i < len(colors) else BLUE
        shape.line.fill.background()

        tf_box = shape.text_frame
        tf_box.word_wrap = True
        tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_box.paragraphs[0].text = step_title
        tf_box.paragraphs[0].font.size = Pt(16)
        tf_box.paragraphs[0].font.bold = True
        tf_box.paragraphs[0].font.color.rgb = WHITE

        if step_desc:
            p_desc = tf_box.add_paragraph()
            p_desc.text = step_desc
            p_desc.font.size = Pt(11)
            p_desc.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            p_desc.alignment = PP_ALIGN.CENTER

        # 矢印テキスト
        if i < len(steps) - 1:
            arrow_box = slide.shapes.add_textbox(
                x_start + box_width // 2 - Inches(0.3),
                y + box_height,
                Inches(0.6), Inches(0.5)
            )
            atf = arrow_box.text_frame
            atf.paragraphs[0].text = "▼"
            atf.paragraphs[0].font.size = Pt(20)
            atf.paragraphs[0].font.color.rgb = GRAY
            atf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


# ══════════════════════════════════════════════════════════
# PPTX 1: 抽出ルール・プロンプト説明資料
# ══════════════════════════════════════════════════════════
def create_rules_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # -- 表紙 --
    add_title_slide(prs,
        "fnl-builder 抽出ルール・プロンプト説明資料",
        "現場担当者との合意形成のために\n生成日: 2026-03-03")

    # -- 目次 --
    add_content_slide(prs, "目次", [
        ("1. 全体像 ─ 抽出品質を担保する処理フロー", Pt(18), False, BLACK, 0),
        ("2. LLM System Prompt の要約（必須ルール・禁則・優先順位）", Pt(18), False, BLACK, 0),
        ("3. LLM Extract Prompt の要約（入力・schema要求・出力制約）", Pt(18), False, BLACK, 0),
        ("4. taxonomy の要約（category/phase の一覧と選び方）", Pt(18), False, BLACK, 0),
        ("5. コース別サプリメントの仕組み", Pt(18), False, BLACK, 0),
        ("6. ルールベース候補抽出の要約", Pt(18), False, BLACK, 0),
        ("7. banned フィルタの要約", Pt(18), False, BLACK, 0),
        ("8. 現場確認ポイント（運用で決めるべき点）", Pt(18), False, BLACK, 0),
    ])

    # -- 1. 全体像 --
    add_section_slide(prs, "1", "全体像 ─ 抽出品質を担保する処理フロー")

    add_flow_slide(prs, "3段階の品質担保パイプライン", [
        ("第1段階：ルールベース候補抽出（広い網）",
         "正規表現キーワード群で「特記事項になりやすい行」を候補として拾い上げる"),
        ("第2段階：LLM による構造化抽出（精密な判定）",
         "taxonomy に沿って category/phase/severity 等に分類し、JSON schema 準拠データを返す"),
        ("第3段階：統合・フィルタ・リライト（最終仕上げ）",
         "banned パターンによる除外、社内メモ除去、LLM出力とルールベース出力の統合"),
    ], [GREEN, PURPLE, ORANGE])

    # -- 2. System Prompt --
    add_section_slide(prs, "2", "LLM System Prompt の要約")

    add_table_slide(prs, "最重要ルール 14項目（前半）",
        ["No.", "ルール名", "内容"],
        [
            ["1", "推測禁止", "入力に書かれていない事実を作らない。不確実なら explicitness=unclear で表現"],
            ["2", "JSON出力のみ", "返信は {\"items\":[...]} のみ。Markdown/説明文は禁止"],
            ["3", "スキーマ厳守", "taxonomy の id のみ使用。ML items schema に準拠"],
            ["4", "Evidence完全一致", "evidence.quote は入力本文からの完全一致コピー"],
            ["5", "caution条件", "severity=caution なら caution_reason 必須"],
            ["6", "PII最小化", "旅券番号、電話番号、メール等を who/details/summary に書かない"],
            ["7", "抽出判定の順番", "「手配会社は誰に何をすればよいか」が明確でなければ抽出しない"],
        ])

    add_table_slide(prs, "最重要ルール 14項目（後半）",
        ["No.", "ルール名", "内容"],
        [
            ["8", "禁則", "金銭/保険/旅券/連絡先変更/座席希望/WT関連等を問答無用で除外"],
            ["9", "航空会社関連", "機内食RQ/搭乗支援のみ採用。窓側/通路側は除外"],
            ["10", "正規化", "SSR/WCHRコード禁止（平文で記述）。新婚旅行→ハネムーンに統一"],
            ["11", "推測抑制(同室)", "同室相手が明示された場合のみ抽出"],
            ["12", "ランドオンリー正規化", "L/O等→「ランドオンリー（航空券自己手配・現地合流）」に統一"],
            ["13", "イレギュラー伝達事項", "定型外でも現地OPが行動すべき具体的指示は採用"],
            ["14", "振替予約・コース混在", "別コースセクションは medical 系のみ採用"],
        ])

    add_content_slide(prs, "禁則の詳細範囲（ルール8）", [
        ("System Prompt で定義される禁則一覧:", Pt(18), True, BLUE, 0),
        ("金銭/請求/入金/領収書/クレカ/残金案内", Pt(15), False, BLACK, 0),
        ("旅行保険", Pt(15), False, BLACK, 0),
        ("旅券（パスポート番号・発行/有効期限・取得予定等）", Pt(15), False, BLACK, 0),
        ("連絡先変更（住所/電話/メール/送付先変更）", Pt(15), False, BLACK, 0),
        ("航空機の座席希望/座席クラス希望", Pt(15), False, BLACK, 0),
        ("旅行手続書類の送付希望 / 各地発着申込", Pt(15), False, BLACK, 0),
        ("時間帯指定なし / 未確定・予定段階のリクエスト", Pt(15), False, BLACK, 0),
        ("マイページ単独通知 / WT（ウェイティング）関連", Pt(15), False, BLACK, 0),
        ("営業/社内担当者向けの進行メモ・連絡済み報告", Pt(15), False, BLACK, 0),
        ("対象者不明のアイテム（who_id を推測で埋めない）", Pt(15), False, BLACK, 0),
    ])

    add_content_slide(prs, "強制抽出ルール & 医療詳細ルール", [
        ("■ 強制抽出ルール", Pt(20), True, GREEN, 0),
        ("アレルギー：否定を除き、必ず items 化（category: dietary, phase: meal_time）", Pt(14), False, BLACK, 0),
        ("クレーム/苦情：否定を除き、vip_sensitive として items 化", Pt(14), False, BLACK, 0),
        ("FNL指示文言：対外アクションが明示される場合のみ items 化", Pt(14), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("■ 医療関連の詳細要約ルール", Pt(20), True, RED, 0),
        ("病名・症状に加え、医療機器/治療器具/薬品/量/本数/サイズ/持込方法を具体記載", Pt(14), False, BLACK, 0),
        ("社内進行メモ（お伺い書待ち/HC登録済み等）は含めない", Pt(14), False, BLACK, 0),
        ("医療事実・数値・薬品名・機器スペックは省略禁止", Pt(14), False, BLACK, 0),
        ("240文字超の場合は同一人物でも items を分割する", Pt(14), False, BLACK, 0),
    ])

    # -- 3. Extract Prompt --
    add_section_slide(prs, "3", "LLM Extract Prompt の要約")

    add_content_slide(prs, "Extract Prompt の入力構造と出力制約", [
        ("■ 入力構造（テンプレートの2つのプレースホルダ）", Pt(20), True, BLUE, 0),
        ("{{TAXONOMY_YAML}} ─ taxonomy.yaml の全文（category/phase の正本）", Pt(15), False, BLACK, 0),
        ("{{PAGES_TEXT}} ─ PDF/CSVから抽出したページ別テキスト（1始まり）", Pt(15), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("■ 出力制約", Pt(20), True, BLUE, 0),
        ("JSON object のみを返す（説明文/Markdown禁止）", Pt(15), False, BLACK, 0),
        ("全ページを最後まで確認する（後半ページの見落とし禁止）", Pt(15), False, BLACK, 0),
        ("該当0件なら {\"items\": []} を返す。無理に作らない", Pt(15), False, BLACK, 0),
        ("handoff_text は運用向けの1文、SSR/WCHRコード禁止", Pt(15), False, BLACK, 0),
    ])

    add_table_slide(prs, "出力 schema（各 item の必須キー）",
        ["キー", "型", "説明"],
        [
            ["category", "string", "taxonomy の categories.id"],
            ["phase", "string", "taxonomy の phases.id"],
            ["summary", "string", "人が即判断できる短文"],
            ["explicitness", "enum", "explicit / implicit / unclear"],
            ["confidence", "float", "0.0〜1.0"],
            ["severity", "enum", "error / warning / caution"],
            ["evidence", "object", "{ page: 整数, quote: 完全一致抜粋 }"],
            ["handoff_text", "string", "運用向けの1文"],
            ["who_id", "string", "対象者の問い合わせ番号-枝番"],
        ])

    # -- 4. taxonomy --
    add_section_slide(prs, "4", "taxonomy の要約")

    add_table_slide(prs, "category 一覧（11カテゴリ）",
        ["id", "ラベル", "対象"],
        [
            ["medical_health", "医療/健康", "病気、怪我、持病、妊娠、服薬"],
            ["dietary", "食事制限/アレルギー", "アレルギー、宗教/嗜好、ベジタリアン"],
            ["mobility_accessibility", "移動支援/バリアフリー", "車椅子、杖、歩行補助、階段不可"],
            ["accommodation_room", "宿泊/部屋/ホテル要望", "禁煙/喫煙、ベッドタイプ、階層"],
            ["grouping_companion", "同行/グルーピング", "同行者、同室/別室希望"],
            ["vip_sensitive", "VIP/要配慮", "VIP、クレームリスク、特別贈呈"],
            ["schedule_change_separation", "離団/日程変更", "離団、途中参加、延泊"],
            ["documents_immigration", "書類/入出国", "ビザ、ESTA、入国条件"],
            ["communication_language", "連絡/言語", "通訳要望、緊急連絡先"],
            ["baggage_equipment", "手荷物/特殊装備", "大型荷物、医療機器"],
            ["other", "その他", "上記に当てはまらない特記事項"],
        ])

    add_table_slide(prs, "phase 一覧（11フェーズ）",
        ["id", "ラベル", "対象場面"],
        [
            ["pre_departure", "出発前", "事前準備、手続"],
            ["departure_airport", "出発空港/集合", "空港集合、出国前"],
            ["flight", "機内/航空移動", "搭乗中、座席配慮"],
            ["arrival_airport", "到着空港/入国", "入国、到着導線"],
            ["transfer", "送迎/移動", "バス、鉄道、車移動"],
            ["on_tour", "現地催行中", "観光地、ガイド対応"],
            ["hotel_stay", "ホテル滞在", "チェックイン、部屋"],
            ["meal_time", "食事", "レストラン、食事制限対応"],
            ["free_time_optional", "自由行動/OP", "オプション"],
            ["return_trip", "帰国/解散", "復路、帰国時"],
            ["unknown", "不明/全般", "フェーズ判断不可"],
        ])

    # -- 5. コース別サプリメント --
    add_section_slide(prs, "5", "コース別サプリメントの仕組み")

    add_content_slide(prs, "コース番号の抽出とマージ手順", [
        ("■ コース番号の抽出ロジック", Pt(20), True, BLUE, 0),
        ("コースコードからアルファベットプレフィックスを除去し数字列を抽出", Pt(15), False, BLACK, 0),
        ("E417 → 417  /  EH417 → 417  /  ET470 → 470", Pt(15), False, PURPLE, 0),
        ("", Pt(10), False, BLACK, 0),
        ("■ マージ手順", Pt(20), True, BLUE, 0),
        ("1. 入力のコースコード群から重複なしの数字列セットを作る", Pt(15), False, BLACK, 0),
        ("2. 各数字に対して courses/{number}.md を検索", Pt(15), False, BLACK, 0),
        ("3. 見つかったファイルを昇順に並べ、重複排除しながらマージ", Pt(15), False, BLACK, 0),
        ("4. 1つも見つからなかった場合は courses/_default.md にフォールバック", Pt(15), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("■ 現存するコースファイル", Pt(20), True, BLUE, 0),
        ("417.md : シェンゲン協定圏のビザ要件重点抽出、複数都市移動の離団・合流注意", Pt(15), False, BLACK, 0),
        ("_default.md : 「コース固有の追加指示はありません。共通ルールに従ってください。」", Pt(15), False, BLACK, 0),
    ])

    # -- 6. ルールベース候補抽出 --
    add_section_slide(prs, "6", "ルールベース候補抽出の要約")

    add_table_slide(prs, "キーワード群の分類",
        ["領域", "代表的なキーワード・パターン"],
        [
            ["食事/アレルギー", "アレルギー、ベジタリアン、ハラール、苦手…食"],
            ["医療/病気", "糖尿病、インシュリン、ペースメーカー、人工関節、透析"],
            ["障害/移動支援", "身障者、障害者、要介護、閉所恐怖、高所恐怖"],
            ["歩行補助", "車椅子、杖、松葉杖、歩行器、歩行…困難/制限"],
            ["宿泊設備", "エレベーター、ベッド…サイズ、ツイン/ダブル…希望"],
            ["VIP/クレーム", "VIP、重要顧客、クレーム、苦情、トラブル"],
            ["離団/合流", "離団、途中参加、途中離脱、途中合流"],
            ["ランドオンリー", "ランドオンリー、LAND ONLY、L/O"],
        ])

    add_content_slide(prs, "候補抽出の流れと目的", [
        ("■ 候補抽出の4ステップ", Pt(20), True, GREEN, 0),
        ("1. 各行に30パターン以上のキーワードマッチ判定", Pt(15), False, BLACK, 0),
        ("2. マッチした行に banned フィルタを適用→該当すれば除外", Pt(15), False, BLACK, 0),
        ("3. PDF日本語スペーシング修復、不要プレフィックス除去", Pt(15), False, BLACK, 0),
        ("4. クリーニング後3文字超なら候補として返却", Pt(15), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("■ 目的と位置づけ", Pt(20), True, GREEN, 0),
        ("「取りこぼしを減らす」ための広い網", Pt(16), True, BLACK, 0),
        ("精度よりも再現率を優先する設計", Pt(15), False, BLACK, 0),
        ("誤検出（偽陽性）は後段の LLM 構造化抽出で除外される", Pt(15), False, BLACK, 0),
    ])

    # -- 7. banned フィルタ --
    add_section_slide(prs, "7", "banned フィルタの要約")

    add_two_column_slide(prs,
        "banned フィルタ：パイプライン側 vs System Prompt 側",
        "パイプライン側 banned（3パターン）",
        [
            "請求/入金/領収書/残金/クレジット/支払/料金",
            "旅行保険/保険",
            "社内進行/社内手配/社内",
            "",
            "→ 最低限の早期除外に留めている",
        ],
        "System Prompt のみに存在する禁則",
        [
            "旅券（パスポート番号等）",
            "連絡先変更（住所/電話/メール等）",
            "航空機の座席希望",
            "旅行手続書類の送付希望",
            "各地発着申込",
            "未確定・予定段階のリクエスト",
            "マイページ単独通知",
            "WT関連 / 営業進行メモ",
        ],
        RED, ORANGE
    )

    # -- 8. 現場確認ポイント --
    add_section_slide(prs, "8", "現場確認ポイント")

    add_content_slide(prs, "運用で決めるべき7つの論点", [
        ("8-1. 禁則の範囲をどこに寄せるか", Pt(17), True, BLUE, 0),
        ("   → パイプライン側を拡張するか、LLM判断に委ねるか", Pt(14), False, GRAY, 0),
        ("8-2. RQ/HK の採用基準", Pt(17), True, BLUE, 0),
        ("   → OP のHK（予約確定済み）でも現地に伝達すべきケースがないか", Pt(14), False, GRAY, 0),
        ("8-3. 別コース混在時の medical 限定採用", Pt(17), True, BLUE, 0),
        ("   → VIP/センシティブ情報も引き継ぐべきか", Pt(14), False, GRAY, 0),
        ("8-4. evidence.quote 完全一致の厳密度", Pt(17), True, BLUE, 0),
        ("   → PDF崩れで正当な項目が不採用になるリスク", Pt(14), False, GRAY, 0),
        ("8-5. 社内メモ除去の範囲", Pt(17), True, BLUE, 0),
        ("   → 新たなメモ表現の追加手順を定めておくべき", Pt(14), False, GRAY, 0),
        ("8-6. 医療系の信頼度閾値（現状: 0.80）", Pt(17), True, BLUE, 0),
        ("   → 閾値を下げると採用増だがノイズも増える", Pt(14), False, GRAY, 0),
        ("8-7. FNL指示文言の対外アクション判定", Pt(17), True, BLUE, 0),
        ("   → 判定基準が LLM 解釈に依存する点の対策", Pt(14), False, GRAY, 0),
    ])

    path = OUT_DIR / "fnl-rules-and-prompts.pptx"
    prs.save(str(path))
    print(f"[OK] {path}")


# ══════════════════════════════════════════════════════════
# PPTX 2: Web UI 操作手順書（スクリーンショット付き）
# ══════════════════════════════════════════════════════════
def _add_screenshot_slide(prs, step_label, title, bullets, img_path,
                          img_on_right=True, caption=""):
    """テキスト＋スクリーンショットの2カラムスライドを追加する.

    img_on_right=True  → 左テキスト / 右スクショ
    img_on_right=False → 上テキスト / 下スクショ（横長画像向け）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SLIDE_BG)

    # ── ステップラベル（左上バッジ） ──
    badge = slide.shapes.add_shape(5, Inches(0.4), Inches(0.25), Inches(1.6), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].text = step_label
    btf.paragraphs[0].font.size = Pt(14)
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if img_on_right:
        # ── 左カラム: タイトル + 箇条書き ──
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(5.6), Inches(6.2))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(12)

        for item in bullets:
            bp = tf.add_paragraph()
            if isinstance(item, tuple):
                text, size, bold, color = item
            else:
                text, size, bold, color = item, Pt(14), False, BLACK
            bp.text = text
            bp.font.size = size
            bp.font.bold = bold
            bp.font.color.rgb = color
            bp.space_after = Pt(4)

        # ── 右カラム: スクリーンショット (白背景カード風) ──
        card = slide.shapes.add_shape(
            5, Inches(6.3), Inches(0.8), Inches(6.7), Inches(5.9)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        card.line.width = Pt(1)
        card.shadow.inherit = False

        if Path(img_path).exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(6.5), Inches(1.0), Inches(6.3), Inches(5.2)
            )
    else:
        # ── 上部: タイトル + 箇条書き ──
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

        for item in bullets:
            bp = tf.add_paragraph()
            if isinstance(item, tuple):
                text, size, bold, color = item
            else:
                text, size, bold, color = item, Pt(14), False, BLACK
            bp.text = text
            bp.font.size = size
            bp.font.bold = bold
            bp.font.color.rgb = color
            bp.space_after = Pt(3)

        # ── 下部: スクリーンショット（横長・中央配置） ──
        card = slide.shapes.add_shape(
            5, Inches(0.8), Inches(3.1), Inches(11.7), Inches(4.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        card.line.width = Pt(1)

        if Path(img_path).exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(1.0), Inches(3.3), Inches(11.3), Inches(3.8)
            )

    # ── キャプション ──
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(6.3), Inches(6.85), Inches(6.7), Inches(0.4))
        ctf = cap_box.text_frame
        ctf.paragraphs[0].text = caption
        ctf.paragraphs[0].font.size = Pt(10)
        ctf.paragraphs[0].font.color.rgb = GRAY
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def create_tutorial_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ss_dir = OUT_DIR / "screenshots"

    # ── 表紙 ──
    add_title_slide(prs,
        "fnl-builder Web UI 操作手順書",
        "現場担当者向け 操作レクチャー資料\n作成日: 2026-03-03")

    # ── 事前準備 ──
    add_content_slide(prs, "事前準備 ─ 用意するファイル", [
        ("操作を開始する前に、以下の3種類のファイルを手元に用意してください。", Pt(16), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("1. RoomingList（PDF）", Pt(18), True, BLUE, 0),
        ("   部屋割り表（ルームタイプ、ゲスト名、問い合わせ番号等）", Pt(14), False, GRAY, 0),
        ("   例: 1027_E417_ROOMINGLIST.pdf", Pt(14), False, PURPLE, 0),
        ("", Pt(8), False, BLACK, 0),
        ("2. PassengerList（PDF）", Pt(18), True, BLUE, 0),
        ("   乗客名簿（氏名、パスポート情報等）", Pt(14), False, GRAY, 0),
        ("   例: 1027_E417_PASSENGERLIST.pdf", Pt(14), False, PURPLE, 0),
        ("", Pt(8), False, BLACK, 0),
        ("3. MessageList（PDF または CSV）", Pt(18), True, BLUE, 0),
        ("   メッセージリスト（特記事項の原文）", Pt(14), False, GRAY, 0),
        ("   例: 1027_E417_MessageList(CSV).csv", Pt(14), False, PURPLE, 0),
    ])

    # ── 全体フロー ──
    add_flow_slide(prs, "操作の全体フロー", [
        ("Step 1: Web UIを開く", "ブラウザでlocalhostにアクセス"),
        ("Step 2-4: 3種類のファイルをアップロード", "RoomingList → PassengerList → MessageList"),
        ("Step 5: LLM Providerを確認", "通常は openai のまま"),
        ("Step 6: 「実行」ボタンを押す", "parse → integrate → render の3段階処理"),
        ("Step 7-8: 結果テーブル＆警告を確認", "名簿テーブルとバリデーション警告を目視チェック"),
        ("Step 9: Excel をダウンロード", "final_list-*.xlsx が自動ダウンロード"),
    ], [BLUE, GREEN, GRAY, RED, ORANGE, PURPLE])

    # ── Step 1: 初期画面 ──
    _add_screenshot_slide(prs, "Step 1", "Web UIを開く", [
        "ブラウザのアドレスバーに URL を入力してアクセスします。",
        "",
        ("左側のサイドバーに4つの入力エリアが表示されます:", Pt(15), True, BLACK),
        "  RoomingList PDF",
        "  PassengerList PDF",
        "  MessageList PDF/CSV",
        "  LLM Provider（ドロップダウン）",
        "",
        ("一番下に赤い「実行」ボタンがあります。", Pt(15), True, RED),
    ], ss_dir / "step1_initial.png",
    caption="初期画面 ─ サイドバーにファイル入力エリアが並ぶ")

    # ── Step 2: Rooming 選択 ──
    _add_screenshot_slide(prs, "Step 2", "RoomingList PDF をアップロード", [
        "「RoomingList PDF」セクションの",
        ("「Browse files」ボタンをクリック", Pt(15), True, BLUE),
        "",
        "ファイル選択ダイアログが開きます。",
        "用意した RoomingList PDF を選択してください。",
        "",
        ("例: 1027_E417_ROOMINGLIST.pdf", Pt(13), False, PURPLE),
    ], ss_dir / "step2_select_rooming.png",
    caption="Browse files でファイル選択ダイアログが開く")

    # ── Step 2 結果: アップロード完了 ──
    _add_screenshot_slide(prs, "Step 2", "RoomingList アップロード完了", [
        "アップロードが完了すると、",
        ("ファイル名とサイズが表示されます。", Pt(15), True, GREEN),
        "",
        "表示を確認したら次のファイルに進みます。",
    ], ss_dir / "step3_rooming_uploaded.png",
    caption="アップロード完了 ─ ファイル名が表示される")

    # ── Step 3: Passenger 選択 ──
    _add_screenshot_slide(prs, "Step 3", "PassengerList PDF をアップロード", [
        "「PassengerList PDF」セクションの",
        ("「Browse files」ボタンをクリック", Pt(15), True, BLUE),
        "",
        "乗客名簿の PDF を選択してください。",
        "",
        ("例: 1027_E417_PASSENGERLIST.pdf", Pt(13), False, PURPLE),
    ], ss_dir / "step4_select_passenger.png",
    caption="PassengerList のファイル選択")

    # ── Step 3 結果 ──
    _add_screenshot_slide(prs, "Step 3", "PassengerList アップロード完了", [
        "アップロードが完了すると、",
        ("ファイル名とサイズが表示されます。", Pt(15), True, GREEN),
    ], ss_dir / "step5_passenger_uploaded.png",
    caption="PassengerList アップロード完了")

    # ── Step 4: MessageList 選択 ──
    _add_screenshot_slide(prs, "Step 4", "MessageList PDF/CSV をアップロード", [
        "「MessageList PDF/CSV」セクションの",
        ("「Browse files」ボタンをクリック", Pt(15), True, BLUE),
        "",
        "メッセージリストのファイルを選択してください。",
        "PDF でも CSV でも対応しています。",
        "",
        ("例: 1027_E417_MessageList(CSV).csv", Pt(13), False, PURPLE),
    ], ss_dir / "step6_select_messagelist.png",
    caption="MessageList のファイル選択（PDF/CSV 両対応）")

    # ── 3ファイル全てアップロード完了 ──
    _add_screenshot_slide(prs, "Step 2-4", "3種類すべてのファイルをアップロード完了", [
        ("3つのファイルがすべてアップロードされた状態です。", Pt(15), True, GREEN),
        "",
        "各セクションにファイル名が表示されていることを確認してください。",
        "",
        ("【重要】", Pt(15), True, ORANGE),
        "3種類すべてのファイルをアップロードしてから",
        "「実行」を押してください。",
    ], ss_dir / "step7_all_uploaded.png",
    caption="全ファイルアップロード完了状態")

    # ── Step 5-6: 実行 ──
    _add_screenshot_slide(prs, "Step 5-6", "LLM Providerを確認して「実行」", [
        ("LLM Provider:", Pt(16), True, BLUE),
        "  サイドバー最下部のドロップダウンを確認",
        "  通常は openai が選択済み。変更不要。",
        "",
        ("赤い「実行」ボタンをクリック", Pt(16), True, RED),
        "",
        "処理開始後の3工程:",
        "  1. parse_stage: PDF/CSV解析",
        "  2. integrate_stage: LLM + ルールベース統合",
        "  3. render_stage: Excel出力",
        "",
        ("【注意】処理中はタブを閉じたり更新しないでください", Pt(14), True, ORANGE),
    ], ss_dir / "step8_click_run.png",
    caption="「実行」ボタンで処理開始")

    # ── Step 7: 結果テーブル ──
    _add_screenshot_slide(prs, "Step 7", "結果テーブルを確認", [
        "処理完了後、メイン画面にツアー情報と",
        ("名簿テーブルが表示されます。", Pt(15), True, BLUE),
        "",
        "表示カラム:",
        "  room_type / number / inquiry",
        "  family_name / given_name / remarks",
        "",
        ("remarks カラムが最終Excelに", Pt(14), True, BLACK),
        ("反映される特記事項です。", Pt(14), True, BLACK),
    ], ss_dir / "step9_result_table.png",
    caption="結果テーブル ─ remarks が特記事項")

    # ── Step 8: バリデーション ──
    _add_screenshot_slide(prs, "Step 8", "バリデーション警告を確認", [
        "テーブル下部に警告が表示される場合があります:",
        "",
        ("  [missing_guest_data]", Pt(14), True, ORANGE),
        "    → パスポート情報が未入力のゲストがいる",
        "",
        ("  [rooms_mismatch_total]", Pt(14), True, ORANGE),
        "    → 部屋数合計が申告数と不一致",
        "",
        ("  [rooms_mismatch_by_type]", Pt(14), True, ORANGE),
        "    → 部屋タイプ別の数が不一致",
        "",
        ("警告が出ても処理は完了。", Pt(14), True, RED),
        ("出力Excelは生成されます。内容を必ず目視チェック。", Pt(14), True, RED),
    ], ss_dir / "step10_validation.png",
    caption="バリデーション警告の例")

    # ── 処理中画面 ──
    _add_screenshot_slide(prs, "参考", "処理中の画面", [
        "「実行」ボタン押下後、",
        "処理の進行状況が表示されます。",
        "",
        ("処理には数十秒〜数分かかる場合があります。", Pt(14), True, BLACK),
        "",
        "進行バーやステータス表示で",
        "現在の処理段階を確認できます。",
    ], ss_dir / "step11_processing.png",
    caption="処理中画面 ─ 進行状況の表示")

    # ── Step 9: ダウンロード ──
    _add_screenshot_slide(prs, "Step 9", "Excel ファイルをダウンロード", [
        "処理完了後、ダウンロード通知が表示されます。",
        "",
        ("ダウンロードされるファイル:", Pt(16), True, BLUE),
        ("  final_list-*.xlsx", Pt(15), True, PURPLE),
        "  （最終名簿 Excel）",
        "",
        "※末尾の数字は連番です。",
        "",
        ("以上で操作は完了です！", Pt(18), True, GREEN),
        "Excelを開き、特記事項の内容を確認してください。",
    ], ss_dir / "step12_download.png",
    caption="完了 ─ Excel ファイルのダウンロード")

    # ── ファイル命名規則 ──
    add_table_slide(prs, "入力ファイルの命名規則",
        ["種別", "命名パターン", "例"],
        [
            ["RoomingList", "{受付番号}_{コース}_ROOMINGLIST.pdf", "1027_E417_ROOMINGLIST.pdf"],
            ["PassengerList", "{受付番号}_{コース}_PASSENGERLIST.pdf", "1027_E417_PASSENGERLIST.pdf"],
            ["MessageList (CSV)", "{受付番号}_{コース}_MessageList(CSV).csv", "1027_E417_MessageList(CSV).csv"],
            ["MessageList (PDF)", "{受付番号}_{コース}_MessageList.pdf", "1027_E417_MessageList.pdf"],
        ])

    # ── トラブルシューティング ──
    add_table_slide(prs, "よくある問題と対処法",
        ["症状", "原因", "対処法"],
        [
            ["「実行」後に何も表示されない", "LLM API接続エラー/ファイル形式不一致", "LLM Provider設定確認。F12でエラー確認"],
            ["remarksが空欄のゲストが多い", "MessageListに特記事項なし/who_id紐付け失敗", "問い合わせ番号-枝番がRoomingListと一致か確認"],
            ["[missing_guest_data]が大量", "PassengerList PDF解析失敗", "より鮮明なPDFを用意する"],
            ["[rooms_mismatch]が出る", "RoomingListの部屋数とPAX数の不整合", "RoomingList原本の部屋数・タイプを確認"],
            ["Excelがダウンロードされない", "処理中エラー/ポップアップブロック", "ダウンロード設定確認。再度「実行」"],
        ])

    path = OUT_DIR / "fnl-builder-tutorial.pptx"
    prs.save(str(path))
    print(f"[OK] {path}")


# ══════════════════════════════════════════════════════════
# PPTX 3: なぜ「ルール＋AI」の2段構成が有効なのか
# ══════════════════════════════════════════════════════════
def create_why_rules_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 表紙
    add_title_slide(prs,
        "なぜ「ルール＋AI」の2段構成が有効なのか",
        "PDF・CSVからの情報抽出において\nコードベースのルール処理とAIを組み合わせることで\n精度・安定性・コストすべてが改善できる理由を解説します",
        DARK_BLUE)

    # 1. AIだけの課題
    add_section_slide(prs, "1", "AIだけで抽出する場合の課題", RED)

    add_content_slide(prs, "AIだけの場合に起きる5つの問題", [
        ("1. 抜け漏れ（見落とし）", Pt(20), True, RED, 0),
        ("   同じPDFでも、実行するたびに抽出される項目が変わることがある", Pt(15), False, BLACK, 0),
        ("", Pt(6), False, BLACK, 0),
        ("2. ばらつき", Pt(20), True, RED, 0),
        ("   「アレルギー」を拾えたり拾えなかったり、結果が不安定", Pt(15), False, BLACK, 0),
        ("", Pt(6), False, BLACK, 0),
        ("3. コスト高", Pt(20), True, RED, 0),
        ("   大量のPDFテキストをすべてAIに送ると、API利用料がかさむ", Pt(15), False, BLACK, 0),
        ("", Pt(6), False, BLACK, 0),
        ("4. 処理時間", Pt(20), True, RED, 0),
        ("   AIの応答を待つため、1件あたり数十秒かかることも", Pt(15), False, BLACK, 0),
        ("", Pt(6), False, BLACK, 0),
        ("5. 原因究明が困難", Pt(20), True, RED, 0),
        ("   なぜ抜け漏れが起きたか後から追跡しにくい", Pt(15), False, BLACK, 0),
    ], RED)

    add_content_slide(prs, "たとえ話：書類チェックの場面で考えると…", [
        ("", Pt(8), False, BLACK, 0),
        ("ベテラン社員に「この100枚の書類から", Pt(20), False, BLACK, 0),
        ("重要な情報を全部拾って」と丸投げするようなもの。", Pt(20), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("いくら優秀でも、", Pt(20), False, BLACK, 0),
        ("チェックリストなしに100枚を読み通せば、", Pt(20), False, BLACK, 0),
        ("当然見落としが出ます。", Pt(20), False, BLACK, 0),
        ("", Pt(14), False, BLACK, 0),
        ("毎回同じ精度は期待できません。", Pt(24), True, RED, 0),
    ], PURPLE)

    # 2. 2段構成とは
    add_section_slide(prs, "2", "「ルール → AI」の2段構成とは", BLUE)

    add_flow_slide(prs, "fnl-builder の2ステップ抽出フロー", [
        ("PDF / CSV ファイル", "ルーミングリスト・乗客名簿・メッセージリスト"),
        ("第1段階：ルールベース抽出", "キーワード照合・正規表現パターン・テンプレート構造解析"),
        ("第2段階：AI による補完抽出", "ルールで取りこぼした情報をAIが文脈で拾う"),
        ("最終出力（Excel）", "ルール抽出 ＋ AI抽出 を統合した完成リスト"),
    ], [BLUE, GREEN, PURPLE, ORANGE])

    add_content_slide(prs, "たとえ話：2段構成＝チェックリスト＋ベテラン", [
        ("第1段階（ルール）= 定型チェックリストで機械的にチェック", Pt(20), True, GREEN, 0),
        ("", Pt(6), False, BLACK, 0),
        ("「アレルギー」「車椅子」「VIP」などの決まったキーワードが", Pt(16), False, BLACK, 0),
        ("書かれていれば 100%確実に拾う。", Pt(16), False, BLACK, 0),
        ("", Pt(14), False, BLACK, 0),
        ("第2段階（AI）= チェックリストにない『行間』をベテランが補足", Pt(20), True, PURPLE, 0),
        ("", Pt(6), False, BLACK, 0),
        ("「足が不自由なので低層階希望」のように、", Pt(16), False, BLACK, 0),
        ("キーワードには当てはまらないが意味的に重要な情報をAIが拾う。", Pt(16), False, BLACK, 0),
    ], BLUE)

    # 3. 役割分担
    add_section_slide(prs, "3", "ルールとAIの役割分担", GREEN)

    add_two_column_slide(prs,
        "それぞれが担当する処理",
        "第1段階：ルールベース抽出",
        [
            "問い合わせ番号の特定（決まった形式を自動認識）",
            "30以上の重要キーワードを一瞬でチェック",
            "  - アレルギー、ベジタリアン、ハラール",
            "  - 糖尿病、ペースメーカー、車椅子",
            "  - 禁煙、低層階、エレベーター近く",
            "  - VIP、苦情、重要顧客",
            "不要情報の除外（保険、社内メモ、経理）",
            "テキスト整形（PDF特有の崩れを修復）",
        ],
        "第2段階：AI補完抽出",
        [
            "文脈の理解",
            "  「足が悪い」→ 移動介助が必要",
            "曖昧な表現の解釈",
            "  「前回ご不満があった」→ VIP・要注意",
            "抜け漏れページの再チェック",
            "  1回目で未抽出のページだけ再スキャン",
            "カテゴリ分類",
            "  「医療」「食事」「移動」に自動分類",
        ],
        GREEN, PURPLE
    )

    add_content_slide(prs, "具体例で見る役割分担", [
        ("【ルールで拾える】", Pt(18), True, GREEN, 0),
        ("原文: 「甲殻類アレルギーのため食事配慮希望」", Pt(15), False, BLACK, 0),
        ("→ 「アレルギー」というキーワードに一致 → 確実に抽出", Pt(14), False, GRAY, 0),
        ("", Pt(8), False, BLACK, 0),
        ("【AIで拾う】", Pt(18), True, PURPLE, 0),
        ("原文: 「母の足が最近弱くなり長距離歩行は避けたい」", Pt(15), False, BLACK, 0),
        ("→ キーワード該当なし → 文脈から「移動介助」と判断", Pt(14), False, GRAY, 0),
        ("", Pt(8), False, BLACK, 0),
        ("【ルールで除外】", Pt(18), True, RED, 0),
        ("原文: 「旅行保険加入済み（証券番号: TI-2024-XXXX）」", Pt(15), False, BLACK, 0),
        ("→ 保険・経理情報は自動フィルタで除外 → AIに渡さない", Pt(14), False, GRAY, 0),
        ("", Pt(8), False, BLACK, 0),
        ("【AIで拾う】", Pt(18), True, PURPLE, 0),
        ("原文: 「結婚25周年の記念旅行とのこと」", Pt(15), False, BLACK, 0),
        ("→ 「記念日」カテゴリとして特別対応の情報に分類", Pt(14), False, GRAY, 0),
    ])

    # 4. 比較
    add_section_slide(prs, "4", "「AIのみ」と「ルール＋AI」の比較", ORANGE)

    add_table_slide(prs, "8項目で比較：AIのみ vs ルール＋AI",
        ["比較項目", "AIのみ", "ルール＋AI"],
        [
            ["定型キーワードの抽出精度", "90%前後（見落としあり）", "100%（ルールで確実に拾う）"],
            ["結果の安定性", "実行ごとにバラつく", "ルール部分は毎回同じ結果"],
            ["全体の抽出率（推定）", "70〜85%", "90%以上"],
            ["処理速度", "全文AI処理 → 遅い", "ルール部分は一瞬 → 全体高速"],
            ["API利用コスト", "全文送信 → 高い", "補完分だけ送信 → 低い"],
            ["不要情報の混入", "AIが誤って拾うことがある", "ルールで事前に除外済み"],
            ["原因追跡", "AIの判断で不透明", "ルール=明確 / AI=根拠付き"],
            ["ルール追加の容易さ", "プロンプト書換え（副作用リスク）", "キーワード1行追加で済む"],
        ])

    # 5. 数字で見る効果
    add_section_slide(prs, "5", "数字で見る効果", BLUE)

    add_content_slide(prs, "100件の重要情報を処理した場合（イメージ）", [
        ("", Pt(8), False, BLACK, 0),
        ("AIのみで抽出できる件数", Pt(18), True, RED, 0),
        ("   70〜85件（15〜30件の漏れ）", Pt(22), True, RED, 0),
        ("", Pt(14), False, BLACK, 0),
        ("ルールだけで確実に抽出できる件数", Pt(18), True, BLUE, 0),
        ("   75件（定型キーワード一致分）", Pt(22), True, BLUE, 0),
        ("", Pt(14), False, BLACK, 0),
        ("ルール＋AIで抽出できる件数", Pt(18), True, GREEN, 0),
        ("   90件以上（漏れが10件以下に）", Pt(22), True, GREEN, 0),
        ("", Pt(14), False, BLACK, 0),
        ("ポイント: 75件分のルール結果は毎回100%同じ → 安定性が格段に向上", Pt(18), True, BLACK, 0),
    ])

    # 6. なぜルールが先なのか
    add_section_slide(prs, "6", "なぜルールを「先に」入れるべきなのか", GREEN)

    add_content_slide(prs, "5つの理由", [
        ("1. 確実に拾えるものは、確実な方法で拾う", Pt(20), True, BLUE, 0),
        ("   キーワード一致で100%確実に、一瞬で抽出", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("2. AIの負担を減らすことで精度が上がる", Pt(20), True, BLUE, 0),
        ("   情報量が少なくなると残りに集中でき、判断精度が逆に上がる", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("3. 不要な情報をAIに見せない", Pt(20), True, BLUE, 0),
        ("   保険情報や社内メモを事前除外→AIの誤判断リスクをなくす", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("4. コスト削減", Pt(20), True, BLUE, 0),
        ("   ルールで先に処理した分はAPIを使わない→コストが下がる", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("5. 問題発生時に原因を特定しやすい", Pt(20), True, BLUE, 0),
        ("   ルール部分は一目瞭然、AI部分も根拠テキスト付き", Pt(15), False, BLACK, 0),
    ])

    # 7. 実際の処理の流れ
    add_section_slide(prs, "7", "実際のシステムではどう動いているか", PURPLE)

    add_content_slide(prs, "4ステップの処理フロー（コード知識不要）", [
        ("ステップ1：ファイルを読み込む", Pt(20), True, BLUE, 0),
        ("   PDF/CSV からテキストを取り出す。CSVは文字コードを自動判定。", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("ステップ2：ルールベース抽出（ここで7〜8割が確定）", Pt(20), True, GREEN, 0),
        ("   30以上のキーワード照合 → 除外パターン適用 → テキスト整形", Pt(15), False, BLACK, 0),
        ("   処理時間は数百ミリ秒（1秒未満）。AIは一切使わない。", Pt(15), True, GREEN, 0),
        ("", Pt(8), False, BLACK, 0),
        ("ステップ3：AI補完抽出（行間を読む）", Pt(20), True, PURPLE, 0),
        ("   整形・フィルタ済みのテキストをAIに渡す → 文脈で拾う", Pt(15), False, BLACK, 0),
        ("   未抽出ページだけ2回目の再スキャンを実施", Pt(15), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("ステップ4：統合 → 最終出力", Pt(20), True, ORANGE, 0),
        ("   ルール結果（土台）＋ AI結果（追加）→ 重複排除 → Excel出力", Pt(15), False, BLACK, 0),
    ])

    # 8. 適用提案
    add_section_slide(prs, "8", "自社AIプロジェクトへの適用提案", DARK_BLUE)

    add_content_slide(prs, "段階的な導入ステップ", [
        ("Step 1: 頻出キーワードのリストアップ", Pt(20), True, BLUE, 0),
        ("   過去の抽出結果を分析し、毎回拾うべき「決まったキーワード」をリスト化", Pt(15), False, BLACK, 0),
        ("   例：アレルギー、車椅子、VIP、禁煙、記念日 など", Pt(14), False, GRAY, 0),
        ("", Pt(10), False, BLACK, 0),
        ("Step 2: 除外パターンの定義", Pt(20), True, BLUE, 0),
        ("   AIに渡す必要がない情報（経理、保険、社内メモ）のパターンを定義", Pt(15), False, BLACK, 0),
        ("   これだけでもAIの精度は向上する", Pt(14), False, GREEN, 0),
        ("", Pt(10), False, BLACK, 0),
        ("Step 3: ルール抽出 → AI補完の2段パイプラインを構築", Pt(20), True, BLUE, 0),
        ("   ルールで拾えたものは「確定情報」→ AIに「残りの見落としを確認」と依頼", Pt(15), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("Step 4: キーワードの継続的追加", Pt(20), True, BLUE, 0),
        ("   AIが新たに拾った頻出情報をルールに昇格 → 網羅率が徐々に向上", Pt(15), False, BLACK, 0),
    ])

    add_content_slide(prs, "たとえ話：業務マニュアルの改善サイクル", [
        ("", Pt(10), False, BLACK, 0),
        ("最初は「気づいたら報告」（＝ AIだけ）", Pt(20), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("↓", Pt(24), False, GRAY, 0),
        ("", Pt(10), False, BLACK, 0),
        ("よくある事例をマニュアル化（＝ ルール化）", Pt(20), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("↓", Pt(24), False, GRAY, 0),
        ("", Pt(10), False, BLACK, 0),
        ("マニュアルに載ってない例外だけ上司判断（＝ AI）", Pt(20), False, BLACK, 0),
        ("", Pt(10), False, BLACK, 0),
        ("↓", Pt(24), False, GRAY, 0),
        ("", Pt(10), False, BLACK, 0),
        ("例外が繰り返されたらマニュアルに追加", Pt(20), False, BLACK, 0),
        ("", Pt(14), False, BLACK, 0),
        ("これが「ルール＋AI」の改善サイクルです。", Pt(24), True, PURPLE, 0),
    ], PURPLE)

    # まとめ
    add_title_slide(prs, "まとめ",
        "1. AIだけに頼ると、抜け漏れ・バラつき・コスト高が同時に起きる\n"
        "2. ルールベースの前処理で「確実に拾えるもの」は100%確実に拾える\n"
        "3. AIは「ルールで拾えなかった部分」だけに集中 → 精度UP・コストDOWN\n"
        "4. ルール＋AIで抽出率を 70-85% → 90%以上に改善可能\n"
        "5. 導入は段階的に可能：キーワードリスト作成から始められる",
        DARK_BLUE)

    # FAQ
    add_content_slide(prs, "よくある質問（FAQ）", [
        ("Q. ルールの作成にはプログラミングの知識が必要？", Pt(18), True, BLUE, 0),
        ("   → ルールの設計（キーワード選定）に知識は不要。コード反映は1行追加で済む", Pt(14), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("Q. AIを使わず、ルールだけではダメ？", Pt(18), True, BLUE, 0),
        ("   → 7〜8割はカバーできるが、文脈依存の情報（残り2〜3割）はAIが必要", Pt(14), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("Q. 導入コストはどれくらい？", Pt(18), True, BLUE, 0),
        ("   → 既存AI処理の「前」にルール追加する形。ランニングコストはむしろ下がる", Pt(14), False, BLACK, 0),
        ("", Pt(8), False, BLACK, 0),
        ("Q. 新しいキーワードが出てきたらどうする？", Pt(18), True, BLUE, 0),
        ("   → AIが新たに拾った頻出情報をルールに追加。使うほどシステムが賢くなる", Pt(14), False, BLACK, 0),
    ])

    path = OUT_DIR / "why-rules-before-ai.pptx"
    prs.save(str(path))
    print(f"[OK] {path}")


# ──────────── メイン ────────────
if __name__ == "__main__":
    create_rules_pptx()
    create_tutorial_pptx()
    create_why_rules_pptx()
    print("\nAll 3 PPTX files generated successfully!")
